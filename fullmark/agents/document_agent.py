"""
fullmark/agents/document_agent.py
----------------------------------
Handles: PDF, DOCX, RTF, TXT, EPUB, XLSX, CSV, ODS, PPTX, ODP, IPYNB, MSG, EML
"""

from __future__ import annotations

import csv
import io
import json
import logging
from datetime import datetime, timezone
from pathlib import Path

from fullmark import AgentError
from fullmark.utils.markdown_utils import (
    clean_text,
    fenced_code,
    front_matter,
    heading,
    rows_to_gfm_table,
)

logger = logging.getLogger(__name__)

_AGENT_NAME = "DocumentAgent"


class DocumentAgent:
    """
    Convert document formats to Markdown.

    Supports: PDF, DOCX, DOC, RTF, TXT, EPUB, XLSX, XLS, CSV, ODS,
              PPTX, PPT, ODP, IPYNB, MSG, EML.
    """

    def convert(self, source: str | Path) -> str:
        """
        Convert *source* document to a Markdown string.

        Args:
            source: Path to the document file.

        Returns:
            Markdown string with YAML front matter.

        Raises:
            AgentError: If the file cannot be read or no converter exists.
        """
        path = Path(source)
        if not path.exists():
            raise AgentError(f"File not found: {path}")

        ext = path.suffix.lower()
        dispatch = {
            ".pdf":  self._convert_pdf,
            ".docx": self._convert_docx,
            ".doc":  self._convert_docx,
            ".rtf":  self._convert_rtf,
            ".txt":  self._convert_txt,
            ".epub": self._convert_epub,
            ".xlsx": self._convert_xlsx,
            ".xls":  self._convert_xlsx,
            ".ods":  self._convert_xlsx,
            ".csv":  self._convert_csv,
            ".pptx": self._convert_pptx,
            ".ppt":  self._convert_pptx,
            ".odp":  self._convert_pptx,
            ".ipynb": self._convert_ipynb,
            ".msg":  self._convert_msg,
            ".eml":  self._convert_eml,
        }

        converter = dispatch.get(ext)
        if converter is None:
            raise AgentError(f"DocumentAgent: unsupported extension '{ext}'")

        logger.debug("converting %s with %s", path.name, converter.__name__)
        body = converter(path)

        # Prepend a ## format section heading if body doesn't already start with one
        format_labels: dict[str, str] = {
            ".pdf":  "PDF Document",
            ".docx": "Word Document",
            ".doc":  "Word Document",
            ".rtf":  "RTF Document",
            ".txt":  "Plain Text",
            ".epub": "EPUB Book",
            ".xlsx": "Spreadsheet",
            ".xls":  "Spreadsheet",
            ".ods":  "Spreadsheet",
            ".csv":  "CSV Data",
            ".pptx": "Slideshow",
            ".ppt":  "Slideshow",
            ".odp":  "Slideshow",
            ".ipynb": "Jupyter Notebook",
            ".msg":  "Email (MSG)",
            ".eml":  "Email (EML)",
        }
        label = format_labels.get(ext, ext.lstrip(".").upper())
        if not body.lstrip().startswith("#"):
            body = f"## {label}: {path.name}\n\n{body}"

        fm = front_matter(path.name, _AGENT_NAME)
        return f"{fm}\n\n{body}"

    # ──────────────────────────────────────────────────────────────────────────
    # PDF
    # ──────────────────────────────────────────────────────────────────────────

    def _convert_pdf(self, path: Path) -> str:
        text = self._pdf_pdfplumber(path)
        if not text.strip():
            logger.debug("pdfplumber returned empty text for %s — trying pdfminer", path.name)
            text = self._pdf_pdfminer(path)
        if not text.strip():
            logger.debug("pdfminer returned empty text for %s — trying tesseract", path.name)
            text = self._pdf_ocr(path)
        return clean_text(text) if text.strip() else f"*Could not extract text from {path.name}*"

    def _pdf_pdfplumber(self, path: Path) -> str:
        try:
            import pdfplumber  # type: ignore
        except ImportError:
            logger.debug("pdfplumber not installed")
            return ""
        parts: list[str] = []
        try:
            with pdfplumber.open(path) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text() or ""
                    tables = page.extract_tables() or []
                    if page_text:
                        parts.append(page_text)
                    for table in tables:
                        if not table:
                            continue
                        headers = [str(c or "") for c in table[0]]
                        rows    = [[str(c or "") for c in row] for row in table[1:]]
                        parts.append(rows_to_gfm_table(headers, rows))
        except Exception as exc:
            logger.warning("pdfplumber failed on %s: %s", path.name, exc)
        return "\n\n".join(parts)

    def _pdf_pdfminer(self, path: Path) -> str:
        try:
            from pdfminer.high_level import extract_text  # type: ignore
        except ImportError:
            logger.debug("pdfminer.six not installed")
            return ""
        try:
            return extract_text(str(path)) or ""
        except Exception as exc:
            logger.warning("pdfminer failed on %s: %s", path.name, exc)
            return ""

    def _pdf_ocr(self, path: Path) -> str:
        try:
            import pytesseract  # type: ignore
            from PIL import Image  # type: ignore
            import fitz  # PyMuPDF — optional for rendering PDF pages  # type: ignore
        except ImportError:
            logger.warning("pytesseract/Pillow/PyMuPDF not available — cannot OCR PDF")
            return ""
        parts: list[str] = []
        try:
            doc = fitz.open(str(path))
            for page in doc:
                pix = page.get_pixmap(dpi=200)
                img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
                text = pytesseract.image_to_string(img)
                if text.strip():
                    parts.append(text)
        except Exception as exc:
            logger.warning("PDF OCR failed on %s: %s", path.name, exc)
        return "\n\n".join(parts)

    # ──────────────────────────────────────────────────────────────────────────
    # DOCX
    # ──────────────────────────────────────────────────────────────────────────

    def _convert_docx(self, path: Path) -> str:
        try:
            from docx import Document  # type: ignore
        except ImportError:
            raise AgentError("python-docx not installed — cannot read DOCX")
        try:
            doc = Document(str(path))
        except Exception as exc:
            # python-docx cannot read legacy binary .doc (OLE format)
            if path.suffix.lower() == ".doc":
                logger.warning("Cannot open %s with python-docx (legacy .doc format). "
                               "Convert to .docx with LibreOffice or Word first.", path.name)
                # Try antiword system binary as a best-effort fallback
                text = self._doc_antiword(path)
                if text.strip():
                    return clean_text(text)
                return (
                    f"## Legacy .doc: {path.name}\n\n"
                    f"> **Note:** This file is in legacy binary `.doc` format which cannot be "
                    f"read directly. Convert to `.docx` format to extract content.\n\n"
                    f"*No text extracted.*"
                )
            raise AgentError(f"Cannot open DOCX {path.name}: {exc}") from exc

        parts: list[str] = []
        for para in doc.paragraphs:
            style = para.style.name if para.style else ""
            text  = para.text.strip()
            if not text:
                continue
            if style.startswith("Heading 1"):
                parts.append(heading(text, 1))
            elif style.startswith("Heading 2"):
                parts.append(heading(text, 2))
            elif style.startswith("Heading 3"):
                parts.append(heading(text, 3))
            elif style.startswith("Heading"):
                parts.append(heading(text, 4))
            else:
                parts.append(text)

        for table in doc.tables:
            headers = [cell.text.strip() for cell in table.rows[0].cells] if table.rows else []
            rows    = [
                [cell.text.strip() for cell in row.cells]
                for row in table.rows[1:]
            ]
            parts.append(rows_to_gfm_table(headers, rows))

        return "\n\n".join(parts)

    # ──────────────────────────────────────────────────────────────────────────
    # RTF
    # ──────────────────────────────────────────────────────────────────────────

    def _doc_antiword(self, path: Path) -> str:
        """Try to extract text from a legacy .doc using the `antiword` system binary."""
        import shutil
        import subprocess
        if not shutil.which("antiword"):
            return ""
        try:
            result = subprocess.run(
                ["antiword", str(path)], capture_output=True, timeout=30
            )
            if result.returncode == 0:
                return result.stdout.decode(errors="replace")
        except Exception as exc:
            logger.debug("antiword failed: %s", exc)
        return ""

    def _convert_rtf(self, path: Path) -> str:
        try:
            from striprtf.striprtf import rtf_to_text  # type: ignore
        except ImportError:
            raise AgentError("striprtf not installed — cannot read RTF")
        raw = path.read_text(encoding="utf-8", errors="replace")
        return clean_text(rtf_to_text(raw))

    # ──────────────────────────────────────────────────────────────────────────
    # TXT
    # ──────────────────────────────────────────────────────────────────────────

    def _convert_txt(self, path: Path) -> str:
        text = path.read_text(encoding="utf-8", errors="replace")
        return clean_text(text)

    # ──────────────────────────────────────────────────────────────────────────
    # EPUB
    # ──────────────────────────────────────────────────────────────────────────

    def _convert_epub(self, path: Path) -> str:
        try:
            import ebooklib  # type: ignore
            from ebooklib import epub  # type: ignore
            from bs4 import BeautifulSoup  # type: ignore
        except ImportError:
            raise AgentError("ebooklib or beautifulsoup4 not installed — cannot read EPUB")

        parts: list[str] = []
        try:
            book = epub.read_epub(str(path))
        except Exception as exc:
            raise AgentError(f"Cannot open EPUB {path.name}: {exc}") from exc

        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_content(), "html.parser")
            for tag in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6"]):
                level = int(tag.name[1])
                parts.append(heading(tag.get_text(strip=True), level))
            for p in soup.find_all("p"):
                text = p.get_text(strip=True)
                if text:
                    parts.append(text)

        return "\n\n".join(parts)

    # ──────────────────────────────────────────────────────────────────────────
    # XLSX / ODS
    # ──────────────────────────────────────────────────────────────────────────

    def _convert_xlsx(self, path: Path) -> str:
        ext = path.suffix.lower()

        # Legacy .xls — try xlrd first, openpyxl can't read it
        if ext == ".xls":
            text = self._xls_xlrd(path)
            if text:
                return text
            return (
                f"## Legacy Spreadsheet: {path.name}\n\n"
                f"> **Note:** `.xls` format requires the `xlrd` package "
                f"(`pip install xlrd`). Install it to extract content.\n\n"
                f"*No data extracted.*"
            )

        try:
            import openpyxl  # type: ignore
        except ImportError:
            raise AgentError("openpyxl not installed — cannot read XLSX/ODS")

        try:
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        except Exception as exc:
            raise AgentError(f"Cannot open spreadsheet {path.name}: {exc}") from exc

        parts: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(heading(f"Sheet: {sheet_name}", 2))
            rows_data = list(ws.iter_rows(values_only=True))
            if not rows_data:
                continue
            headers = [str(c) if c is not None else "" for c in rows_data[0]]
            rows    = [
                [str(c) if c is not None else "" for c in row]
                for row in rows_data[1:]
            ]
            parts.append(rows_to_gfm_table(headers, rows))

        return "\n\n".join(parts)

    # ──────────────────────────────────────────────────────────────────────────
    # CSV
    # ──────────────────────────────────────────────────────────────────────────

    def _xls_xlrd(self, path: Path) -> str:
        """Try to read a legacy .xls file using xlrd. Returns empty string if unavailable."""
        try:
            import xlrd  # type: ignore
        except ImportError:
            logger.debug("xlrd not installed — cannot read .xls (pip install xlrd)")
            return ""
        try:
            wb = xlrd.open_workbook(str(path))
        except Exception as exc:
            logger.warning("xlrd failed on %s: %s", path.name, exc)
            return ""

        parts: list[str] = []
        for sheet in wb.sheets():
            parts.append(heading(f"Sheet: {sheet.name}", 2))
            if sheet.nrows == 0:
                continue
            headers = [str(sheet.cell_value(0, c)) for c in range(sheet.ncols)]
            rows: list[list[str]] = []
            for r in range(1, sheet.nrows):
                rows.append([str(sheet.cell_value(r, c)) for c in range(sheet.ncols)])
            parts.append(rows_to_gfm_table(headers, rows))
        return "\n\n".join(parts)

    def _convert_csv(self, path: Path) -> str:
        raw = path.read_text(encoding="utf-8", errors="replace")
        reader = csv.reader(io.StringIO(raw))
        rows_data = list(reader)
        if not rows_data:
            return "*Empty CSV file*"
        headers = rows_data[0]
        rows    = rows_data[1:]
        return rows_to_gfm_table(headers, rows)

    # ──────────────────────────────────────────────────────────────────────────
    # PPTX / ODP
    # ──────────────────────────────────────────────────────────────────────────

    def _convert_pptx(self, path: Path) -> str:
        try:
            from pptx import Presentation  # type: ignore
            from pptx.util import Pt  # noqa: F401
        except ImportError:
            raise AgentError("python-pptx not installed — cannot read PPTX")

        try:
            prs = Presentation(str(path))
        except Exception as exc:
            raise AgentError(f"Cannot open PPTX {path.name}: {exc}") from exc

        parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            slide_title = ""
            texts: list[str] = []
            notes_text = ""

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if not line:
                            continue
                        if shape.shape_type == 13 or (hasattr(shape, "placeholder_format")
                                and shape.placeholder_format
                                and shape.placeholder_format.idx == 0):
                            slide_title = line
                        else:
                            texts.append(line)

            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame
                if notes:
                    notes_text = notes.text.strip()

            parts.append(heading(slide_title or f"Slide {i}", 2))
            if texts:
                parts.extend(f"- {t}" for t in texts)
            if notes_text:
                parts.append(f"\n> **Notes:** {notes_text}")

        return "\n\n".join(parts)

    # ──────────────────────────────────────────────────────────────────────────
    # IPYNB (Jupyter Notebook)
    # ──────────────────────────────────────────────────────────────────────────

    def _convert_ipynb(self, path: Path) -> str:
        try:
            nb = json.loads(path.read_text(encoding="utf-8"))
        except Exception as exc:
            raise AgentError(f"Cannot parse IPYNB {path.name}: {exc}") from exc

        parts: list[str] = []
        cells = nb.get("cells", [])
        for cell in cells:
            cell_type = cell.get("cell_type", "")
            source    = "".join(cell.get("source", []))
            if cell_type == "markdown":
                parts.append(source)
            elif cell_type == "code":
                lang = nb.get("metadata", {}).get("kernelspec", {}).get("language", "python")
                parts.append(fenced_code(source, lang))
                # Include text outputs
                for output in cell.get("outputs", []):
                    otype = output.get("output_type", "")
                    if otype in ("stream", "execute_result", "display_data"):
                        out_text = "".join(output.get("text", []) or
                                          output.get("data", {}).get("text/plain", []))
                        if out_text.strip():
                            parts.append(fenced_code(out_text.strip(), ""))
            elif cell_type == "raw":
                parts.append(fenced_code(source))

        return "\n\n".join(parts)

    # ──────────────────────────────────────────────────────────────────────────
    # MSG (Outlook)
    # ──────────────────────────────────────────────────────────────────────────

    def _convert_msg(self, path: Path) -> str:
        try:
            import extract_msg  # type: ignore
        except ImportError:
            raise AgentError("extract-msg not installed — cannot read MSG")

        try:
            msg = extract_msg.Message(str(path))
        except Exception as exc:
            raise AgentError(f"Cannot open MSG {path.name}: {exc}") from exc

        parts = [
            f"**From:** {msg.sender or ''}",
            f"**To:** {msg.to or ''}",
            f"**Subject:** {msg.subject or ''}",
            f"**Date:** {msg.date or ''}",
            "",
            msg.body or "",
        ]
        return "\n\n".join(parts)

    # ──────────────────────────────────────────────────────────────────────────
    # EML (RFC 2822 email)
    # ──────────────────────────────────────────────────────────────────────────

    def _convert_eml(self, path: Path) -> str:
        import email as _email
        from email import policy as _policy

        raw = path.read_bytes()
        msg = _email.message_from_bytes(raw, policy=_policy.default)

        parts = [
            f"**From:** {msg.get('From', '')}",
            f"**To:** {msg.get('To', '')}",
            f"**Subject:** {msg.get('Subject', '')}",
            f"**Date:** {msg.get('Date', '')}",
            "",
        ]

        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                ct = part.get_content_type()
                if ct == "text/plain":
                    body = part.get_content()
                    break
                if ct == "text/html" and not body:
                    try:
                        from bs4 import BeautifulSoup  # type: ignore
                        body = BeautifulSoup(part.get_content(), "html.parser").get_text()
                    except ImportError:
                        body = part.get_content()
        else:
            body = msg.get_content()

        parts.append(clean_text(body or ""))
        return "\n\n".join(parts)
